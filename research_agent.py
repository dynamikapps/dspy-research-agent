import sys
import os
import dspy
import asyncio
import json
import pandas as pd
from rich import print
from typing import List, Dict, Any, Optional, Tuple
from tavily import TavilyClient
from crawl4ai import AsyncWebCrawler, BrowserConfig, CrawlerRunConfig
from dsp.trackers.langfuse_tracker import LangfuseTracker
from dotenv import load_dotenv
from datetime import datetime
from pathlib import Path
from config import config
from tenacity import retry, stop_after_attempt, wait_exponential
import traceback
from io import StringIO

load_dotenv()

# Set up API keys
tavily_api_key = os.getenv('TAVILY_API_KEY')
stability_api_key = os.getenv('STABILITY_API_KEY')  # For image generation

if not all([tavily_api_key]):
    print("Missing required API keys in environment variables.")
    sys.exit(1)

# Initialize clients and trackers
langfuse = LangfuseTracker(trace_name="advanced_research_agent")
tavily_client = TavilyClient(api_key=tavily_api_key)

# Configure DSPy with our config module
dspy.configure(lm=config.get_dspy_lm())

# Create reports directory and subdirectories
REPORTS_DIR = Path("reports")
REPORTS_DIR.mkdir(exist_ok=True)

# Create subdirectory for content
REPORTS_CONTENT_DIR = REPORTS_DIR / "content"
REPORTS_CONTENT_DIR.mkdir(exist_ok=True)

# Create subdirectory for visualizations
REPORTS_VIZ_DIR = REPORTS_DIR / "visualizations"
REPORTS_VIZ_DIR.mkdir(exist_ok=True)


class ResearchPlan(dspy.Signature):
    """Signature for planning research strategy using ReAct."""
    query: str = dspy.InputField()
    context: List[str] = dspy.InputField(
        desc="Initial context from search results")
    research_plan: List[str] = dspy.OutputField(
        desc="List of research steps")
    required_tools: List[str] = dspy.OutputField(desc="List of tools needed")
    follow_up_questions: List[str] = dspy.OutputField(
        desc="Questions to investigate further")


class ResearchOutline(dspy.Signature):
    """Define a signature for generating research outlines based on gathered context."""
    query: str = dspy.InputField()
    context: List[str] = dspy.InputField(
        desc="Research context from search results")
    title: str = dspy.OutputField()
    sections: list[str] = dspy.OutputField()
    section_subheadings: dict[str, list[str]] = dspy.OutputField(
        desc="mapping from section headings to subheadings")
    section_contexts: dict[str, List[int]] = dspy.OutputField(
        desc="mapping from sections to relevant context indices")


class ContentGenerator(dspy.Signature):
    """Signature for generating section content with source attribution."""
    context: List[str] = dspy.InputField()
    metadata: List[Dict] = dspy.InputField()
    title: str = dspy.InputField()
    section_heading: str = dspy.InputField()
    section_subheadings: list[str] = dspy.InputField()
    relevant_context_indices: List[int] = dspy.InputField()
    content: str = dspy.OutputField()
    source_references: List[Dict[str, Any]] = dspy.OutputField(
        desc="List of references with format: {'text': str, 'context_index': int, 'start': int, 'end': int}")
    table_data: List[Dict[str, Any]] = dspy.OutputField(
        desc="List of table data with format: {'data': Any, 'title': str}")


class ContentReviewer(dspy.Signature):
    """Signature for reviewing and improving content."""
    content: str = dspy.InputField()
    context: List[str] = dspy.InputField()
    improvements: str = dspy.OutputField()
    suggestions: List[str] = dspy.OutputField()


class KeyConceptsSignature(dspy.Signature):
    """Signature for extracting key concepts from research context."""
    context: str = dspy.InputField(desc="Research context to analyze")
    concepts: List[str] = dspy.OutputField(
        desc="List of key concepts extracted from the context")


class UserFeedback(dspy.Signature):
    """Signature for processing user feedback and improving content."""
    content: str = dspy.InputField()
    feedback: str = dspy.InputField()
    improved_content: str = dspy.OutputField()
    changes_made: List[str] = dspy.OutputField()


class SectionExpansion(dspy.Signature):
    """Signature for expanding sections based on user feedback."""
    section_content: str = dspy.InputField()
    expansion_request: str = dspy.InputField()
    context: List[str] = dspy.InputField()
    expanded_content: str = dspy.OutputField()
    new_citations: List[Dict[str, Any]] = dspy.OutputField()


# Tool Functions


async def crawl_urls(urls: List[str]) -> List[Dict[str, Any]]:
    """Crawl multiple URLs using Crawl4AI."""
    browser_config = BrowserConfig(headless=True, verbose=True)
    run_config = CrawlerRunConfig(
        word_count_threshold=50,
        excluded_tags=['nav', 'footer', 'header'],
        exclude_external_links=True,
        wait_for="body",
        page_timeout=30000,
        magic=True
    )

    results = []
    async with AsyncWebCrawler(config=browser_config) as crawler:
        for url in urls:
            try:
                result = await crawler.arun(url=url, config=run_config)
                if result.success:
                    results.append({
                        'url': url,
                        'content': result.markdown,
                        'title': result.title
                    })
            except Exception as e:
                print(f"Error crawling {url}: {str(e)}")

    return results


def search_tavily(query: str, max_results: int = 5) -> tuple[List[str], List[Dict[str, str]], List[str]]:
    """Perform a Tavily search and return results."""
    try:
        search_result = tavily_client.search(
            query=query,
            search_depth="advanced",
            max_results=max_results,
            include_answer=True,
            include_raw_content=True
        )

        contexts = []
        metadata = []
        urls_to_crawl = []

        if search_result.get('answer'):
            contexts.append(search_result['answer'])
            metadata.append({
                'title': 'Tavily Synthesized Answer',
                'url': None
            })

        for result in search_result.get('results', []):
            if result.get('content'):
                contexts.append(result['content'])
                metadata.append({
                    'title': result.get('title', 'No title'),
                    'url': result.get('url')
                })
                if result.get('url'):
                    urls_to_crawl.append(result['url'])

        return contexts, metadata, urls_to_crawl

    except Exception as e:
        print(f"Error searching Tavily: {str(e)}")
        return [], [], []


class AdvancedResearchAgent(dspy.Module):
    """Advanced multi-stage pipeline for comprehensive research with ReAct capabilities."""

    def __init__(self):
        super().__init__()
        self.planner = dspy.ChainOfThought(ResearchPlan)
        self.outline_generator = dspy.ChainOfThought(ResearchOutline)
        self.content_generator = dspy.ChainOfThought(ContentGenerator)
        self.content_reviewer = dspy.ChainOfThought(ContentReviewer)
        self.concept_extractor = dspy.ChainOfThought(KeyConceptsSignature)
        self.feedback_processor = dspy.ChainOfThought(UserFeedback)
        self.expander = dspy.ChainOfThought(SectionExpansion)

        # Initialize ReAct for dynamic research actions
        self.react = dspy.ReAct(
            "research_query -> research_findings",
            tools=[
                search_tavily,
                self.generate_table
            ]
        )

    async def gather_research(self, query: str, max_depth: int = 2) -> tuple[List[str], List[Dict]]:
        """Gather research data with specified depth."""
        try:
            all_contexts = []
            all_sources = []
            processed_urls = set()  # Track processed URLs to avoid duplicates

            print("\nPerforming initial research...")
            # Initial search with rate limiting and retries
            try:
                contexts, metadata, urls = await self._rate_limited_search(query)
                if contexts:
                    all_contexts.extend(contexts)
                    all_sources.extend(metadata)
                    # Track processed URLs
                    for source in metadata:
                        if source.get('url'):
                            processed_urls.add(source['url'])
            except Exception as e:
                print(f"Warning: Initial search error: {str(e)}")

            # Crawl initial URLs
            if urls:
                print("Crawling initial URLs...")
                crawled_data = await crawl_urls([url for url in urls[:3] if url not in processed_urls])
                for item in crawled_data:
                    if item['url'] not in processed_urls:
                        all_contexts.append(item['content'])
                        all_sources.append(
                            {'title': item['title'], 'url': item['url']})
                        processed_urls.add(item['url'])

            # Perform deeper research if needed
            if max_depth > 1 and all_contexts:  # Only proceed if we have initial context
                print("\nExtracting key concepts for deeper research...")
                # Extract key concepts from initial research
                combined_context = "\n\n".join(all_contexts[:5])
                key_concepts = self.extract_key_concepts(combined_context)

                if key_concepts:
                    print(
                        f"Identified key concepts: {', '.join(key_concepts)}")

                    # Research each key concept with rate limiting
                    for concept in key_concepts[:3]:
                        print(f"\nResearching concept: {concept}")
                        try:
                            sub_contexts, sub_metadata, sub_urls = await self._rate_limited_search(f"{query} {concept}")
                            if sub_contexts:
                                # Add new contexts and sources
                                all_contexts.extend(sub_contexts)
                                for source in sub_metadata:
                                    if source.get('url') not in processed_urls:
                                        all_sources.append(source)
                                        processed_urls.add(source['url'])

                                # Crawl new URLs
                                new_urls = [url for url in sub_urls[:2]
                                            if url not in processed_urls]
                                if new_urls:
                                    sub_crawled = await crawl_urls(new_urls)
                                    for item in sub_crawled:
                                        if item['url'] not in processed_urls:
                                            all_contexts.append(
                                                item['content'])
                                            all_sources.append(
                                                {'title': item['title'], 'url': item['url']})
                                            processed_urls.add(item['url'])
                        except Exception as e:
                            print(
                                f"Warning: Error researching concept '{concept}': {str(e)}")
                            continue

            print(
                f"\nResearch gathering complete. Found {len(all_contexts)} content pieces from {len(all_sources)} sources.")
            return all_contexts, all_sources

        except Exception as e:
            print(f"Error during research gathering: {str(e)}")
            raise

    async def _rate_limited_search(self, query: str) -> Tuple[List[str], List[Dict], List[str]]:
        """Perform rate-limited Tavily search with retries."""
        @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
        async def _search():
            try:
                return search_tavily(query)
            except Exception as e:
                if "429" in str(e):  # Rate limit error
                    print("Rate limit reached, waiting before retry...")
                    raise  # This will trigger the retry
                elif "400" in str(e):  # Bad request
                    print(f"Invalid search query: {query}")
                    return [], [], []
                else:
                    print(f"Search error: {str(e)}")
                    return [], [], []

        return await _search()

    def extract_key_concepts(self, context: str) -> List[str]:
        """Extract key concepts from research context."""
        try:
            result = self.concept_extractor(context=context)
            return result.concepts if hasattr(result, 'concepts') else []
        except Exception as e:
            print(f"Error extracting key concepts: {str(e)}")
            return []

    async def generate_report(self, query: str, user_feedback_callback=None) -> Dict[str, Any]:
        """Generate comprehensive research report with user feedback integration."""
        try:
            print("\nGathering research data...")
            all_contexts, all_sources = await self.gather_research(query, max_depth=2)

            if not all_contexts:
                raise ValueError(
                    "No research data found. Please try a different query or check your search configuration.")

            print("Planning research approach...")
            plan = self.planner(
                query=query,
                context=all_contexts[:5] if len(
                    all_contexts) >= 5 else all_contexts
            )

            print("Generating detailed outline...")
            outline = self.outline_generator(
                query=query,
                context=all_contexts
            )

            print("Generating content for each section...")
            sections = []
            citations = []
            citation_id = 1

            for heading, subheadings in outline.section_subheadings.items():
                print(f"\nProcessing section: {heading}")

                # Get relevant context indices and validate them
                relevant_indices = outline.section_contexts.get(heading, [])
                valid_indices = [
                    i for i in relevant_indices if 0 <= i < len(all_contexts)]

                if not valid_indices and relevant_indices:
                    print(
                        f"Warning: No valid context indices found for section '{heading}'")
                    # Use first few contexts as fallback
                    relevant_contexts = all_contexts[:3]
                else:
                    relevant_contexts = [all_contexts[i]
                                         for i in valid_indices]

                # Generate initial content with source citations
                content_result = self.content_generator(
                    context=all_contexts,
                    metadata=all_sources,
                    title=outline.title,
                    section_heading=heading,
                    section_subheadings=subheadings,
                    relevant_context_indices=valid_indices
                )

                # Get user feedback if callback is provided
                if user_feedback_callback:
                    feedback = user_feedback_callback({
                        'section': heading,
                        'content': content_result.content,
                        'table_data': content_result.table_data
                    })

                    if feedback:
                        if feedback.get('content_feedback'):
                            improved_content, changes = self.process_user_feedback(
                                content_result.content,
                                feedback['content_feedback']
                            )
                            content_result.content = improved_content

                        if feedback.get('expansion_request'):
                            expanded_content, new_citations = self.expand_section(
                                content_result.content,
                                feedback['expansion_request'],
                                relevant_contexts
                            )
                            content_result.content = expanded_content

                            for citation in new_citations:
                                citation['id'] = citation_id
                                citations.append(citation)
                                citation_id += 1

                # Review and improve content
                review_result = self.content_reviewer(
                    content=content_result.content,
                    context=relevant_contexts
                )

                # Process citations
                section_citations = []
                for citation in content_result.source_references:
                    context_index = citation.get('context_index', -1)
                    if 0 <= context_index < len(all_sources):
                        source = all_sources[context_index]
                        section_citations.append({
                            'id': citation_id,
                            'location': citation['start'],
                            'source': source,
                            'text': citation['text']
                        })
                        citation_id += 1

                # Process tables
                section_tables = []
                for table_data in content_result.table_data:
                    if not isinstance(table_data, dict) or 'data' not in table_data:
                        print(
                            f"Warning: Invalid table data in section '{heading}'")
                        continue

                    if isinstance(table_data['data'], (list, pd.DataFrame)):
                        table_content = self.generate_table(table_data['data'])
                        if table_content:
                            section_tables.append(table_content)

                sections.append({
                    'heading': heading,
                    'subheadings': subheadings,
                    'content': review_result.improvements,
                    'citations': section_citations,
                    'tables': section_tables,
                    'suggestions': review_result.suggestions
                })
                citations.extend(section_citations)

            print("\nFinalizing report...")
            report_data = {
                'title': outline.title,
                'sections': sections,
                'sources': all_sources,
                'citations': citations,
                'outline': {
                    'sections': outline.sections,
                    'section_subheadings': outline.section_subheadings,
                    'section_contexts': outline.section_contexts
                },
                'research_plan': plan.research_plan,
                'follow_up_questions': plan.follow_up_questions
            }

            # Final user feedback on complete report if callback is provided
            if user_feedback_callback:
                final_feedback = user_feedback_callback({
                    'type': 'final_review',
                    'report': report_data
                })

                if final_feedback and final_feedback.get('content_feedback'):
                    for section in report_data['sections']:
                        if section['heading'] in final_feedback['content_feedback']:
                            improved_content, changes = self.process_user_feedback(
                                section['content'],
                                final_feedback['content_feedback'][section['heading']]
                            )
                            section['content'] = improved_content

            return report_data

        except Exception as e:
            print(f"Error generating report: {str(e)}")
            raise

    def generate_table(self, data) -> Optional[str]:
        """Generate a markdown table from data."""
        try:
            if isinstance(data, pd.DataFrame):
                return data.to_markdown(index=False)
            elif isinstance(data, list):
                df = pd.DataFrame(data)
                return df.to_markdown(index=False)
            else:
                print(f"Warning: Unsupported table data type: {type(data)}")
                return None
        except Exception as e:
            print(f"Error generating table: {str(e)}")
            return None

    def process_user_feedback(self, content: str, feedback: str) -> Tuple[str, List[str]]:
        """Process user feedback and improve content."""
        try:
            result = self.feedback_processor(
                content=content, feedback=feedback)
            return result.improved_content, result.changes_made
        except Exception as e:
            print(f"Error processing feedback: {str(e)}")
            return content, []

    def expand_section(self, section_content: str, expansion_request: str, context: List[str]) -> Tuple[str, List[Dict[str, Any]]]:
        """Expand a section based on user feedback."""
        try:
            result = self.expander(
                section_content=section_content,
                expansion_request=expansion_request,
                context=context
            )
            return result.expanded_content, result.new_citations
        except Exception as e:
            print(f"Error expanding section: {str(e)}")
            return section_content, []

    def export_to_docx(self, report: Dict[str, Any], output_path: Path) -> Path:
        """Export the report to a Word document."""
        try:
            from docx import Document
            from docx.shared import Inches

            doc = Document()
            doc.add_heading(report['title'], 0)

            # Add timestamp
            doc.add_paragraph(
                f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")

            # Add research methodology
            doc.add_heading('Research Methodology', 1)
            for i, step in enumerate(report['research_plan'], 1):
                doc.add_paragraph(f"{i}. {step}")

            # Add sections
            for section in report['sections']:
                doc.add_heading(section['heading'], 1)

                # Add content
                doc.add_paragraph(section['content'])

                # Add tables
                if section.get('tables'):
                    doc.add_heading('Tables', 2)
                    for table in section['tables']:
                        # Convert markdown table to docx table
                        try:
                            # Parse markdown table into pandas DataFrame
                            df = pd.read_csv(
                                StringIO(table), sep='|', skipinitialspace=True)
                            # Remove empty columns
                            df = df.dropna(axis=1, how='all')

                            # Create docx table
                            table_rows = len(df) + 1  # +1 for header
                            table_cols = len(df.columns)
                            docx_table = doc.add_table(
                                rows=table_rows, cols=table_cols)
                            docx_table.style = 'Table Grid'

                            # Add headers
                            for j, column in enumerate(df.columns):
                                docx_table.cell(0, j).text = str(
                                    column).strip()

                            # Add data
                            for i, row in enumerate(df.values, start=1):
                                for j, cell in enumerate(row):
                                    docx_table.cell(i, j).text = str(
                                        cell).strip()

                            doc.add_paragraph()  # Add spacing after table
                        except Exception as table_error:
                            print(
                                f"Warning: Could not convert table in section {section['heading']}: {str(table_error)}")
                            # Add as plain text instead
                            doc.add_paragraph(table)

                # Add citations
                if section.get('citations'):
                    doc.add_heading('Sources', 2)
                    for citation in section['citations']:
                        doc.add_paragraph(
                            f"[{citation['id']}] {citation['text']} - {citation['source'].get('url', 'N/A')}",
                            style='List Number'
                        )

            # Add follow-up questions
            if report.get('follow_up_questions'):
                doc.add_heading('Follow-up Questions', 1)
                for question in report['follow_up_questions']:
                    doc.add_paragraph(f"• {question}", style='List Bullet')

            # Save document
            doc_path = output_path.with_suffix('.docx')
            doc.save(doc_path)
            return doc_path

        except Exception as e:
            print(f"Error exporting to Word: {str(e)}")
            return None

    def export_to_pdf(self, report: Dict[str, Any], output_path: Path) -> Path:
        """Export the report to a PDF document."""
        try:
            import pdfkit
            from io import StringIO
            import platform

            # Check if wkhtmltopdf is installed and configure it
            if platform.system() == 'Darwin':  # macOS
                config = pdfkit.configuration(
                    wkhtmltopdf='/usr/local/bin/wkhtmltopdf')
            elif platform.system() == 'Linux':
                config = pdfkit.configuration(
                    wkhtmltopdf='/usr/bin/wkhtmltopdf')
            elif platform.system() == 'Windows':
                config = pdfkit.configuration(
                    wkhtmltopdf=r'C:\Program Files\wkhtmltopdf\bin\wkhtmltopdf.exe')
            else:
                config = None

            # First generate HTML content
            html = f"""
            <html>
            <head>
                <meta charset="UTF-8">
                <title>{report['title']}</title>
                <style>
                    body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                    h1 {{ color: #2c3e50; }}
                    h2 {{ color: #34495e; margin-top: 30px; }}
                    table {{ border-collapse: collapse; width: 100%; margin: 20px 0; }}
                    th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
                    th {{ background-color: #f5f5f5; }}
                    .citation {{ font-size: 0.9em; color: #7f8c8d; }}
                    .methodology {{ background-color: #f9f9f9; padding: 20px; border-radius: 5px; }}
                    pre {{ white-space: pre-wrap; word-wrap: break-word; }}
                </style>
            </head>
            <body>
                <h1>{report['title']}</h1>
                <p><em>Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</em></p>
                
                <div class="methodology">
                <h2>Research Methodology</h2>
                <ol>
            """

            # Add research methodology
            for step in report['research_plan']:
                html += f"<li>{step}</li>"
            html += "</ol></div>"

            # Add sections
            for section in report['sections']:
                html += f"<h2>{section['heading']}</h2>"
                html += f"<div>{section['content']}</div>"

                # Add tables
                if section.get('tables'):
                    html += "<h3>Tables</h3>"
                    for table in section['tables']:
                        try:
                            # Convert markdown table to HTML table
                            df = pd.read_csv(
                                StringIO(table), sep='|', skipinitialspace=True)
                            df = df.dropna(axis=1, how='all')
                            html += df.to_html(index=False, classes='table')
                        except Exception as table_error:
                            print(
                                f"Warning: Could not convert table: {str(table_error)}")
                            html += f"<pre>{table}</pre>"

                # Add citations
                if section.get('citations'):
                    html += '<div class="citation"><h3>Sources:</h3><ol>'
                    for citation in section['citations']:
                        html += f"<li>{citation['text']} - {citation['source'].get('url', 'N/A')}</li>"
                    html += "</ol></div>"

            # Add follow-up questions
            if report.get('follow_up_questions'):
                html += "<h2>Follow-up Questions</h2><ul>"
                for question in report['follow_up_questions']:
                    html += f"<li>{question}</li>"
                html += "</ul>"

            html += "</body></html>"

            # Save HTML temporarily
            html_path = output_path.with_suffix('.html')
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write(html)

            # PDF conversion options
            options = {
                'encoding': 'UTF-8',
                'margin-top': '20mm',
                'margin-right': '20mm',
                'margin-bottom': '20mm',
                'margin-left': '20mm',
                'enable-local-file-access': None,
                'quiet': None
            }

            # Convert to PDF
            pdf_path = output_path.with_suffix('.pdf')
            try:
                pdfkit.from_file(str(html_path), str(pdf_path),
                                 options=options, configuration=config)
            except OSError as e:
                if 'wkhtmltopdf' in str(e):
                    raise Exception(
                        "wkhtmltopdf is not installed. Please install it first:\n"
                        "- On macOS: brew install wkhtmltopdf\n"
                        "- On Linux: sudo apt-get install wkhtmltopdf\n"
                        "- On Windows: Download and install from https://wkhtmltopdf.org/downloads.html"
                    ) from e
                raise

            # Clean up HTML file
            html_path.unlink()

            return pdf_path

        except Exception as e:
            print(f"Error exporting to PDF: {str(e)}")
            raise  # Re-raise the exception to show the detailed error message

    def export_to_pptx(self, report: Dict[str, Any], output_path: Path) -> Path:
        """Export the report to a PowerPoint presentation."""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            prs = Presentation()

            # Title slide
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = report['title']
            title_slide.placeholders[1].text = f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"

            # Content slides
            for section in report['sections']:
                # Section title slide
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = section['heading']

                # Content slides
                content_chunks = self._chunk_content(section['content'])
                for chunk in content_chunks:
                    content_slide = prs.slides.add_slide(prs.slide_layouts[2])
                    content_slide.shapes.title.text = section['heading']
                    content_slide.placeholders[1].text = chunk

                # Visualization slides
                for viz in section.get('visualizations', []):
                    if viz.get('path'):
                        viz_path = output_path.parent / viz['path']
                        if viz_path.exists():
                            viz_slide = prs.slides.add_slide(
                                prs.slide_layouts[6])
                            viz_slide.shapes.title.text = viz.get(
                                'title', 'Visualization')
                            pic = viz_slide.shapes.add_picture(
                                str(viz_path),
                                Inches(1),
                                Inches(1.5),
                                width=Inches(8)
                            )

            # Save presentation
            pptx_path = output_path.with_suffix('.pptx')
            prs.save(pptx_path)
            return pptx_path

        except Exception as e:
            print(f"Error exporting to PowerPoint: {str(e)}")
            return None

    def _convert_to_html(self, report: Dict[str, Any]) -> str:
        """Convert report to HTML format."""
        html = f"""
        <html>
        <head>
            <title>{report['title']}</title>
            <style>
                body {{ font-family: Arial, sans-serif; line-height: 1.6; margin: 40px; }}
                h1 {{ color: #2c3e50; }}
                h2 {{ color: #34495e; margin-top: 30px; }}
                img {{ max-width: 100%; height: auto; margin: 20px 0; }}
                .citation {{ font-size: 0.9em; color: #7f8c8d; }}
            </style>
        </head>
        <body>
            <h1>{report['title']}</h1>
            <p><em>Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}</em></p>
        """

        for section in report['sections']:
            html += f"<h2>{section['heading']}</h2>"
            html += f"<div>{section['content']}</div>"

            # Add visualizations
            for viz in section.get('visualizations', []):
                if viz.get('path'):
                    html += f"<img src='{viz['path']}' alt='{viz.get('title', 'Visualization')}'>"

            # Add citations
            if section.get('citations'):
                html += "<div class='citation'><h3>Sources:</h3><ol>"
                for citation in section['citations']:
                    html += f"<li>{citation['text']} - {citation['source'].get('url', 'N/A')}</li>"
                html += "</ol></div>"

        html += "</body></html>"
        return html

    def _chunk_content(self, content: str, max_chars: int = 1000) -> List[str]:
        """Split content into chunks for PowerPoint slides."""
        words = content.split()
        chunks = []
        current_chunk = []
        current_length = 0

        for word in words:
            if current_length + len(word) + 1 <= max_chars:
                current_chunk.append(word)
                current_length += len(word) + 1
            else:
                chunks.append(' '.join(current_chunk))
                current_chunk = [word]
                current_length = len(word)

        if current_chunk:
            chunks.append(' '.join(current_chunk))

        return chunks


def format_markdown_report(report: Dict[str, Any]) -> str:
    """Format the research report in markdown with internal source linking."""
    try:
        timestamp = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

        md = f"# {report['title']}\n\n"
        md += f"*Generated on: {timestamp}*\n\n"

        # Table of Contents
        md += "## Table of Contents\n\n"
        for section in report['outline']['sections']:
            section_link = section.lower().replace(' ', '-')
            md += f"- [{section}](#{section_link})\n"
            for subheading in report['outline']['section_subheadings'].get(section, []):
                subheading_link = subheading.lower().replace(' ', '-')
                md += f"  - [{subheading}](#{subheading_link})\n"
        md += "\n"

        # Research Plan
        md += "## Research Methodology\n\n"
        for i, step in enumerate(report['research_plan'], 1):
            md += f"{i}. {step}\n"
        md += "\n"

        # Research Outline
        md += "## Research Outline\n\n"
        for section in report['outline']['sections']:
            md += f"- {section}\n"
            for subheading in report['outline']['section_subheadings'].get(section, []):
                md += f"  - {subheading}\n"
        md += "\n"

        # Main Content with Citations
        for section in report['sections']:
            md += f"## {section['heading']}\n\n"

            # Add content with citation processing
            content = section['content']
            citations = sorted(section.get('citations', []),
                               key=lambda x: x.get('location', 0), reverse=True)

            # Process citations from end to start to maintain correct positions
            for citation in citations:
                if citation.get('location') is not None and citation.get('id') is not None:
                    citation_mark = f"[^{citation['id']}]"
                    location = citation['location']
                    if 0 <= location <= len(content):
                        content = content[:location] + \
                            citation_mark + content[location:]

            md += content + "\n\n"

            # Add tables
            if section.get('tables'):
                md += "### Tables\n\n"
                for table in section['tables']:
                    md += f"{table}\n\n"

            # Add citations for this section
            if citations:
                md += "\n**Sources:**\n\n"
                for citation in sorted(citations, key=lambda x: x.get('id', 0)):
                    source = citation.get('source', {})
                    md += f"[^{citation['id']}]: {citation['text']} - {source.get('url', 'N/A')}\n"
                md += "\n"

        # Follow-up Questions
        if report.get('follow_up_questions'):
            md += "## Follow-up Questions\n\n"
            for question in report['follow_up_questions']:
                md += f"- {question}\n"
            md += "\n"

        return md

    except Exception as e:
        print(f"Error formatting markdown report: {str(e)}")
        traceback.print_exc()
        return f"Error generating report: {str(e)}"


def save_report(content: str, query: str, report: Dict[str, Any]) -> Path:
    """Save the report and its assets to the reports directory."""
    try:
        # Sanitize the query for use in filename
        sanitized_query = "".join(
            c if c.isalnum() else "_" for c in query)[:50]
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Create report directory
        report_dir = REPORTS_CONTENT_DIR / f"{sanitized_query}_{timestamp}"
        report_dir.mkdir(parents=True, exist_ok=True)

        # Save the main report
        report_file = report_dir / "report.md"
        with open(report_file, 'w', encoding='utf-8') as f:
            f.write(content)

        # Save metadata
        metadata = {
            'query': query,
            'timestamp': timestamp,
            'title': report['title'],
            'outline': report['outline'],
            'research_plan': report.get('research_plan', []),
            'follow_up_questions': report.get('follow_up_questions', []),
            'sources': report['sources']
        }

        with open(report_dir / "metadata.json", 'w', encoding='utf-8') as f:
            json.dump(metadata, f, indent=2, default=str)

        return report_file
    except Exception as e:
        print(f"Error saving report: {str(e)}")
        raise


async def main():
    print("\n=== Advanced Research Agent ===")
    print(f"Reports will be saved in: {REPORTS_DIR.absolute()}")
    print("Enter your research query. Type 'exit' or 'quit' to end.\n")

    while True:
        query = input("Research query: ").strip()

        if query.lower() in ['exit', 'quit']:
            print("\nThank you for using the Advanced Research Agent. Goodbye!")
            break

        if not query:
            print("Please enter a valid query.")
            continue

        print("\nGenerating comprehensive research report...")
        agent = AdvancedResearchAgent()
        report = await agent.generate_report(query)

        # Generate markdown report
        markdown_report = format_markdown_report(report)

        # Save report to file
        file_path = save_report(markdown_report, query, report)

        print(f"\nReport generated and saved to: {file_path}")
        print(f"Report directory: {file_path.parent}")
        print("\n" + "-"*50 + "\n")

if __name__ == "__main__":
    asyncio.run(main())
